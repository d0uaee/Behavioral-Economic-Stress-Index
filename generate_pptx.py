"""
generate_pptx.py — Génération PowerPoint BESI Maroc
Utilise python-pptx pour créer une présentation professionnelle 18 diapositives
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import copy
from pathlib import Path

ROOT = Path(__file__).resolve().parent
OUT  = ROOT / "outputs"
FIG  = OUT / "figures"

# ── Palette couleurs ──────────────────────────────────────────────────────────
C_NAVY    = RGBColor(0x1A, 0x2E, 0x4A)   # bleu marine foncé (fond titres)
C_BLUE    = RGBColor(0x1F, 0x77, 0xB4)   # bleu principal
C_ORANGE  = RGBColor(0xFF, 0x7F, 0x0E)   # orange BESI
C_GREEN   = RGBColor(0x2C, 0xA0, 0x2C)   # vert succès
C_RED     = RGBColor(0xD6, 0x27, 0x28)   # rouge rejet
C_GREY    = RGBColor(0xF5, 0xF5, 0xF5)   # gris très clair (fond body)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK    = RGBColor(0x22, 0x22, 0x22)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # layout entièrement vide

# ─────────────────────────────────────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_w=Pt(0)):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    shape.line.width = line_w
    if line_rgb:
        shape.line.color.rgb = line_rgb
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             font_size=18, bold=False, color=C_DARK,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def title_slide_bg(slide, title, subtitle="", authors=""):
    """Diapositive de titre (fond marine, accent orange)."""
    # Fond entier
    add_rect(slide, 0, 0, 13.33, 7.5, C_NAVY)
    # Bande orange gauche
    add_rect(slide, 0, 0, 0.25, 7.5, C_ORANGE)
    # Titre principal
    add_text(slide, title,   0.5, 1.8, 12.3, 1.6,
             font_size=32, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    if subtitle:
        add_text(slide, subtitle, 0.5, 3.5, 12.3, 0.8,
                 font_size=20, color=C_ORANGE, align=PP_ALIGN.CENTER)
    if authors:
        add_text(slide, authors, 0.5, 5.8, 12.3, 0.5,
                 font_size=14, color=C_WHITE, align=PP_ALIGN.CENTER, italic=True)


def section_slide(slide, num, title, subtitle=""):
    """Diapositive de section (fond bleu)."""
    add_rect(slide, 0, 0, 13.33, 7.5, C_BLUE)
    add_rect(slide, 0, 0, 0.3,  7.5, C_ORANGE)
    add_text(slide, f"PARTIE {num}", 0.6, 2.5, 12, 0.6,
             font_size=16, color=C_ORANGE, align=PP_ALIGN.CENTER, bold=True)
    add_text(slide, title,         0.6, 3.1, 12, 1.2,
             font_size=30, color=C_WHITE, align=PP_ALIGN.CENTER, bold=True)
    if subtitle:
        add_text(slide, subtitle,  0.6, 4.4, 12, 0.6,
                 font_size=16, color=C_WHITE, align=PP_ALIGN.CENTER, italic=True)


def content_slide(slide, title, body_lines, num_label=""):
    """Diapositive de contenu standard."""
    # Bande titre
    add_rect(slide, 0, 0, 13.33, 1.05, C_NAVY)
    add_rect(slide, 0, 0, 0.2, 1.05, C_ORANGE)
    if num_label:
        add_text(slide, num_label, 0.25, 0.1, 0.8, 0.85,
                 font_size=10, color=C_ORANGE, bold=True)
    add_text(slide, title, 0.35, 0.1, 12.5, 0.85,
             font_size=22, bold=True, color=C_WHITE)
    # Fond body
    add_rect(slide, 0, 1.05, 13.33, 6.45, C_GREY)
    # Corps texte
    body = "\n".join(body_lines)
    add_text(slide, body, 0.4, 1.2, 12.5, 6.1,
             font_size=14, color=C_DARK, wrap=True)


def two_col_slide(slide, title, left_lines, right_lines, num_label=""):
    """Diapositive deux colonnes."""
    add_rect(slide, 0, 0, 13.33, 1.05, C_NAVY)
    add_rect(slide, 0, 0, 0.2, 1.05, C_ORANGE)
    if num_label:
        add_text(slide, num_label, 0.25, 0.1, 0.8, 0.85,
                 font_size=10, color=C_ORANGE, bold=True)
    add_text(slide, title, 0.35, 0.1, 12.5, 0.85,
             font_size=22, bold=True, color=C_WHITE)
    add_rect(slide, 0, 1.05, 13.33, 6.45, C_GREY)
    # Séparateur vertical
    add_rect(slide, 6.6, 1.2, 0.04, 6.1, C_BLUE)
    left  = "\n".join(left_lines)
    right = "\n".join(right_lines)
    add_text(slide, left,  0.4, 1.2, 6.1, 6.1, font_size=13, color=C_DARK)
    add_text(slide, right, 6.8, 1.2, 6.1, 6.1, font_size=13, color=C_DARK)


def metrics_slide(slide, title, rows, headers, num_label="", best_row=1):
    """Diapositive avec tableau de métriques stylisé."""
    add_rect(slide, 0, 0, 13.33, 1.05, C_NAVY)
    add_rect(slide, 0, 0, 0.2, 1.05, C_ORANGE)
    if num_label:
        add_text(slide, num_label, 0.25, 0.1, 0.8, 0.85,
                 font_size=10, color=C_ORANGE, bold=True)
    add_text(slide, title, 0.35, 0.1, 12.5, 0.85,
             font_size=22, bold=True, color=C_WHITE)
    add_rect(slide, 0, 1.05, 13.33, 6.45, C_GREY)

    # En-têtes
    col_w = 12.5 / len(headers)
    for ci, h in enumerate(headers):
        add_rect(slide, 0.4 + ci*col_w, 1.25, col_w-0.05, 0.45, C_BLUE)
        add_text(slide, h, 0.4 + ci*col_w, 1.25, col_w-0.05, 0.45,
                 font_size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    row_h = 0.52
    for ri, row in enumerate(rows):
        t = 1.75 + ri * row_h
        bg = RGBColor(0xFF, 0xF3, 0xCD) if ri == best_row else (
             RGBColor(0xE8, 0xF4, 0xF8) if ri % 2 == 0 else C_WHITE)
        for ci, cell in enumerate(row):
            add_rect(slide, 0.4 + ci*col_w, t, col_w-0.05, row_h-0.04, bg,
                     line_rgb=RGBColor(0xCC, 0xCC, 0xCC), line_w=Pt(0.5))
            bold_cell = (ri == best_row)
            add_text(slide, str(cell), 0.42 + ci*col_w, t+0.04,
                     col_w-0.1, row_h-0.1,
                     font_size=12, bold=bold_cell, color=C_DARK,
                     align=PP_ALIGN.CENTER)


# =============================================================================
# DIAPOSITIVE 1 — TITRE
# =============================================================================
sl = prs.slides.add_slide(BLANK)
title_slide_bg(
    sl,
    "Detection Precoce du Stress Economique des Menages au Maroc",
    "BESI — Behavioral Economic Stress Index",
    "Douae Ahadji & Adama Basse  |  Séries Temporelles — ENSAM Meknes  |  Mai 2026"
)

# =============================================================================
# DIAPOSITIVE 2 — ETAT D'AVANCEMENT
# =============================================================================
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 1.05, C_NAVY)
add_rect(sl, 0, 0, 0.2, 1.05, C_ORANGE)
add_text(sl, "ETAT D'AVANCEMENT DU PROJET", 0.35, 0.1, 12.5, 0.85,
         font_size=22, bold=True, color=C_WHITE)
add_rect(sl, 0, 1.05, 13.33, 6.45, C_GREY)

modules = [
    ("Donnees & Pipeline",    "IPC HCP + Google Trends + BESI composite",        "TERMINE", C_GREEN),
    ("Stationnarité + SARIMA","ADF/KPSS, ACF/PACF, grille AIC, walk-forward",    "TERMINE", C_GREEN),
    ("SARIMAX + BESI",        "BESI_trends, comparaison v2, sous-periodes",       "TERMINE", C_GREEN),
    ("Analyse statistique",   "Chow, Granger, Early Warning, Markov",             "TERMINE", C_GREEN),
    ("Deep Learning",         "LSTM 4 fenetres × 2 configs + Prophet",            "TERMINE", C_GREEN),
    ("NLP Maroc",             "Scraping presse + scoring Darija/Arabe/Francais",  "TERMINE", C_GREEN),
    ("Notebooks",             "4 notebooks avec outputs (exploration→résultats)", "TERMINE", C_GREEN),
    ("Rapport & Presentation","README + PRESENTATION_FINALE + PowerPoint",        "TERMINE", C_GREEN),
]

row_h = 0.68
for i, (mod, desc, status, col) in enumerate(modules):
    t = 1.15 + i * row_h
    bg = RGBColor(0xF0, 0xFA, 0xF0) if i % 2 == 0 else C_WHITE
    add_rect(sl, 0.3, t, 9.5, row_h-0.06, bg,
             line_rgb=RGBColor(0xDD,0xDD,0xDD), line_w=Pt(0.5))
    add_rect(sl, 9.85, t, 3.1, row_h-0.06, col)
    add_text(sl, f"  {mod}", 0.35, t+0.08, 3.2, row_h-0.15,
             font_size=13, bold=True, color=C_NAVY)
    add_text(sl, desc,       3.6,  t+0.08, 6.2, row_h-0.15,
             font_size=12, color=C_DARK)
    add_text(sl, f"✓  {status}", 9.9, t+0.1, 2.9, row_h-0.2,
             font_size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

add_text(sl, "Avancement global :  100 %   —   8 / 8 modules complets",
         0.3, 6.75, 12.5, 0.5,
         font_size=14, bold=True, color=C_BLUE, align=PP_ALIGN.CENTER)

# =============================================================================
# DIAPOSITIVE 3 — CONTEXTE & H1
# =============================================================================
sl = prs.slides.add_slide(BLANK)
two_col_slide(
    sl,
    "Contexte & Hypothese de Recherche",
    [
        "PROBLEMATIQUE",
        "",
        "Les statistiques officielles de l'IPC",
        "(HCP) sont publiees avec 1 mois de",
        "delai et ne capturent pas le ressenti",
        "des menages en temps reel.",
        "",
        "=> Peut-on detecter le stress",
        "   economique AVANT qu'il",
        "   apparaisse dans l'IPC ?",
        "",
        "CONTEXTE",
        "",
        "Choc inflationniste 2022 :",
        "  Inflation YoY max = 6.57% (mars 2022)",
        "  Record depuis 30 ans au Maroc",
        "  Guerre Ukraine => prix energie +50%",
    ],
    [
        "HYPOTHESE H1",
        "",
        "\"Les signaux comportementaux",
        "digitaux (Google Trends, Reddit,",
        "YouTube) integres dans BESI",
        "permettent de detecter le stress",
        "economique 1 a 2 mois avant",
        "l'IPC du HCP.\"",
        "",
        "CONTRIBUTION ORIGINALE",
        "",
        "=> BESI : 1er indice composite",
        "   marocain base sur signaux",
        "   digitaux (Darija+Arabe+Fr)",
        "",
        "=> Pipeline NLP multilingue",
        "=> Comparaison 7 modeles",
    ],
    num_label="02"
)

# =============================================================================
# DIAPOSITIVE 4 — DONNEES & BESI
# =============================================================================
sl = prs.slides.add_slide(BLANK)
content_slide(
    sl,
    "Donnees & Indice BESI",
    [
        "VARIABLE CIBLE : IPC mensuel Maroc (HCP)  |  2010-01 a 2024-12  |  180 observations  |  freq=MS",
        "",
        "SIGNAUX COMPORTEMENTAUX COLLECTES :",
        "  • Google Trends (pytrends, geo=MA) : 7 mots-cles — \"inflation maroc\", \"prix huile\",",
        "    \"hausse prix\", \"credit consommation\", \"chomage maroc\", \"prix alimentaires\", \"pouvoir achat\"",
        "  • Reddit r/Morocco (praw)     : keywords inflation / prix / cherte / economie",
        "  • YouTube Data API v3         : 4 chaines marocaines (2M, Medi1TV, Hespress TV, Medias24)",
        "  • Presse marocaine (NLP)      : Hespress, Le360, Medias24, L'Economiste",
        "",
        "TROIS VERSIONS DE BESI (tous signaux normalises 0-1) :",
        "",
        "  BESI composite  = 0.40 × Trends + 0.30 × Reddit + 0.20 × YouTube + 0.10 × |dIPC|",
        "  BESI_trends     = 0.70 × Trends + 0.30 × |dIPC|            <- VERSION ROBUSTE (donnees reelles uniquement)",
        "  BESI_enrichi    = 0.35 × Trends + 0.25 × NLP_Maroc + 0.20 × YouTube + 0.10 × Reddit + 0.10 × |dIPC|",
        "",
        "BESI_trends : moyenne = 0.376  |  ecart-type = 0.228  |  Pic : choc inflationniste 2022 (confirme)",
    ],
    num_label="03"
)

# =============================================================================
# DIAPOSITIVE 5 — STATIONNARITÉ & SARIMA
# =============================================================================
sl = prs.slides.add_slide(BLANK)
two_col_slide(
    sl,
    "Stationnarite & Identification SARIMA",
    [
        "TESTS ADF + KPSS",
        "",
        "IPC niveau    =>  Non stationnaire",
        "  ADF p > 0.05 | KPSS p < 0.05",
        "",
        "IPC diff(1)   =>  Stationnaire",
        "  ADF p < 0.05 | KPSS p > 0.05",
        "",
        "IPC diff(1).diff(12)  =>  Stationnaire",
        "  => SARIMA avec d=1, D=1, s=12",
        "",
        "DECOMPOSITION STL",
        "",
        "  Tendance : +39% sur 14 ans",
        "  Saisonnalite : stable (±0.01)",
        "  Residu : choc visible 2022",
    ],
    [
        "IDENTIFICATION SARIMA",
        "",
        "ACF  : decroissance lente, pics 12/24",
        "  => saisonnalite s=12, D=1",
        "",
        "PACF : coupure nette au lag 2",
        "  => AR(2), p=2",
        "",
        "ACF residuelle : pic lag 1 seul",
        "  => MA(1), q=1",
        "",
        "MODELE OPTIMAL (grille 6 modeles) :",
        "",
        "  SARIMA(2,1,1) x (0,1,1)[12]",
        "  AIC = -502.56  |  BIC = -492.34",
        "",
        "  Residus : bruit blanc",
        "  Ljung-Box p > 0.05  => OK",
    ],
    num_label="04"
)

# =============================================================================
# DIAPOSITIVE 6 — TABLE COMPARAISON MODELES
# =============================================================================
sl = prs.slides.add_slide(BLANK)
metrics_slide(
    sl,
    "Comparaison des Modeles — Walk-Forward h=1 (Test 2022-2024, n=36 mois)",
    headers=["Modele", "RMSE", "MAE", "MAPE", "AIC", "Gain vs SARIMA"],
    rows=[
        ["Naif (Random Walk)",        "0.00409", "0.00339", "0.28%", "—",       "-50.3%"],
        ["★ SARIMA(2,1,1)x(0,1,1)[12]","0.00272","0.00232", "0.19%", "-502.56", "0.0%"],
        ["SARIMAX + Trends",           "0.00327", "0.00253", "0.21%", "-515.48", "-20.3%"],
        ["SARIMAX + BESI_trends ★★",   "0.00304", "0.00241", "0.20%", "-545.88", "-11.9%"],
        ["LSTM (window=12)",           "0.01885", "0.01647", "1.37%", "—",       "-593%"],
        ["Prophet (multiplicatif)",    "0.06082", "0.05950", "4.88%", "—",       "-2136%"],
    ],
    best_row=1,
    num_label="05"
)
add_text(sl,
    "★ Meilleur RMSE   |   ★★ Meilleur AIC : BESI_trends apporte une information statistique reelle",
    0.3, 6.9, 12.7, 0.45,
    font_size=11, bold=True, color=C_BLUE, align=PP_ALIGN.CENTER)

# =============================================================================
# DIAPOSITIVE 7 — PERFORMANCES PAR SOUS-PERIODE
# =============================================================================
sl = prs.slides.add_slide(BLANK)
metrics_slide(
    sl,
    "Performances par Sous-Periode (RMSE Walk-Forward)",
    headers=["Modele", "Choc 2022 (n=12)", "Post-Choc 2023-24 (n=24)", "Test Complet (n=36)"],
    rows=[
        ["Naif (Random Walk)", "0.00597", "0.00269", "0.00409"],
        ["SARIMA",             "0.00197", "0.00303", "0.00272"],
        ["SARIMAX_T",          "0.00380", "0.00297", "0.00327"],
        ["SARIMAX_BT",         "0.00362", "0.00271", "0.00304"],
    ],
    best_row=1,
    num_label="06"
)
add_text(sl,
    "Cle : SARIMA domine pendant le choc 2022  |  SARIMAX_BT legerement meilleur post-2022",
    0.3, 5.7, 12.7, 0.45,
    font_size=13, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)
add_text(sl,
    [
        "Interpretation : La variable exogene BESI perturbe la prevision en periode de rupture (choc imprévisible).",
        "Une fois le regime stabilise post-2022, BESI apporte une information marginale utile (+0.2% gain RMSE).",
    ][0] + "\n" + [
        "Interpretation : La variable exogene BESI perturbe la prevision en periode de rupture (choc imprévisible).",
        "Une fois le regime stabilise post-2022, BESI apporte une information marginale utile (+0.2% gain RMSE).",
    ][1],
    0.4, 6.15, 12.5, 0.9,
    font_size=12, color=C_DARK
)

# =============================================================================
# DIAPOSITIVE 8 — RUPTURE STRUCTURELLE 2022
# =============================================================================
sl = prs.slides.add_slide(BLANK)
two_col_slide(
    sl,
    "Rupture Structurelle 2022 & Dynamique Markov",
    [
        "TEST DE CHOW",
        "",
        "Breakpoint : 2022-01-01",
        "Resultat   : RUPTURE CONFIRMEE",
        "             p < 0.05",
        "",
        "Coefficients pre/post 2022 :",
        "  Constante : +10%",
        "  Tendance  : +60%",
        "  Coeff BESI: +67%",
        "",
        "=> BESI joue un role plus fort",
        "   en periode de crise",
        "",
        "TEST CUSUM",
        "  Depassement bandes mi-2022",
        "  => Changement de regime permanent",
    ],
    [
        "MATRICE DE TRANSITION MARKOV",
        "",
        "               Normal  Warning  H.Stress",
        "Normal    =>   82.2%   17.8%     0.0%",
        "Warning   =>   19.3%   74.2%     6.5%",
        "H.Stress  =>    2.3%    6.8%    90.9%",
        "",
        "CLE : P(High Stress | High Stress)",
        "                        = 90.9%",
        "",
        "=> Une fois en crise, le stress",
        "   est PERSISTANT (non transitoire)",
        "",
        "=> Le choc 2022 = changement de",
        "   regime durable, non episode",
        "   temporaire",
    ],
    num_label="07"
)

# =============================================================================
# DIAPOSITIVE 9 — GRANGER & EARLY WARNING
# =============================================================================
sl = prs.slides.add_slide(BLANK)
two_col_slide(
    sl,
    "Causalite de Granger & Early Warning",
    [
        "CAUSALITE DE GRANGER",
        "",
        "7 / 7 features significatives (p<0.05) :",
        "",
        "  trends_composite   => p < 0.05",
        "  reddit_composite   => p < 0.05",
        "  youtube_composite  => p < 0.05",
        "  besi               => p < 0.05",
        "  ipc_change         => p < 0.05",
        "  ipc_mom            => p < 0.05",
        "  ipc_yoy            => p < 0.05",
        "",
        "=> BESI Granger-cause l'IPC :",
        "   information predictive reelle,",
        "   non spurieuse",
    ],
    [
        "EARLY WARNING — RESULTAT",
        "",
        "Evenement detecte :",
        "  Onset   : 2021-05-01",
        "  Lead time : 12 mois",
        "  avant pic inflation mars 2022",
        "  Detection : CONFIRMEE (TRUE)",
        "",
        "Performance systeme d'alerte :",
        "",
        "  Rappel    = 100%",
        "  (0 faux negatif, 0 crise ratee)",
        "",
        "  F1-Score  = ~0.82",
        "",
        "=> BESI a sonne l'alerte 12 mois",
        "   avant le pic d'inflation record",
    ],
    num_label="08"
)

# =============================================================================
# DIAPOSITIVE 10 — LSTM
# =============================================================================
sl = prs.slides.add_slide(BLANK)
metrics_slide(
    sl,
    "Deep Learning — LSTM : Comparaison Configurations",
    headers=["Configuration", "RMSE", "MAE", "MAPE", "Epochs"],
    rows=[
        ["LSTM IPC seul (window=12) — base",  "0.01885", "0.01647", "1.37%", "19"],
        ["LSTM + BESI (window=12)",            "0.06334", "0.06219", "5.10%", "9"],
        ["LSTM window=6  (sans exog)",         "0.08842", "0.08751", "7.18%", "13"],
        ["LSTM window=12 (sans exog)",         "0.05467", "0.05438", "4.47%", "11"],
        ["LSTM window=18 (sans exog)",         "0.05743", "0.05710", "4.70%", "11"],
        ["LSTM window=24 (sans exog) — best",  "0.04627", "0.04566", "3.76%", "11"],
    ],
    best_row=5,
    num_label="09"
)
add_text(sl,
    "Conclusion : LSTM (window=24) = meilleure config LSTM mais RESTE 17x pire que SARIMA (0.04627 vs 0.00272)\n"
    "Ajouter BESI au LSTM DEGRADE les performances => BESI est un signal 12 mois, incompatible avec LSTM court terme",
    0.4, 6.45, 12.5, 0.9,
    font_size=12, bold=False, color=C_RED)

# =============================================================================
# DIAPOSITIVE 11 — PROPHET
# =============================================================================
sl = prs.slides.add_slide(BLANK)
two_col_slide(
    sl,
    "Prophet — Modele Bayesien de Prevision",
    [
        "CONFIGURATION",
        "",
        "yearly_seasonality = True",
        "weekly_seasonality = False",
        "daily_seasonality  = False",
        "seasonality_mode   = multiplicative",
        "",
        "Train : 2010-01 -> 2021-12",
        "        (144 observations)",
        "Test  : 2022-01 -> 2024-12",
        "        (36 observations)",
        "",
        "Installation automatique lors",
        "de l'execution de run_v2.py",
    ],
    [
        "RESULTATS PROPHET",
        "",
        "  RMSE  =  0.06082",
        "  MAE   =  0.05950",
        "  MAPE  =  4.88%",
        "",
        "=> 22x pire que SARIMA",
        "",
        "POURQUOI PROPHET ECHOUE ?",
        "",
        "Prophet est concu pour :",
        "  - Longues series (annees de jours)",
        "  - Saisonnalite forte et reguliere",
        "  - Sans rupture structurelle majeure",
        "",
        "IPC mensuel + rupture 2022",
        "= contre-indication Prophet",
        "",
        "=> Resultat scientifiquement valide",
    ],
    num_label="10"
)

# =============================================================================
# DIAPOSITIVE 12 — NLP
# =============================================================================
sl = prs.slides.add_slide(BLANK)
content_slide(
    sl,
    "NLP Presse Marocaine — Pipeline Multilingue",
    [
        "MODULE 1 — Scraping presse : Hespress, Le360, Medias24, L'Economiste (selecteurs CSS multi-fallback)",
        "MODULE 2 — Commentaires YouTube : 4 chaines (2M Maroc, Medi1TV, Hespress TV, Medias24)",
        "MODULE 3 — Scoring NLP :",
        "  score = (0.6 × keyword_score + 0.4 × intensite) × engagement_weight  →  normalise 0-1",
        "",
        "DICTIONNAIRE 80+ MOTS-CLES (Darija / Arabe / Francais) :",
        "",
        "  Prix eleves  : ghali, ghla, cher | غالي, غلاء, ارتفاع الأسعار",
        "  Manque argent: ma b9ach, flouss, mskine | فلوس, فقر, محتاج",
        "  Frustration  : hshuma, crise | عيب, أزمة",
        "  Produits base: zit, sokkar, carburant | زيت, سكر, وقود",
        "",
        "MODULE 4 — Aggregation mensuelle (pondere par nb commentaires + likes + vues)",
        "MODULE 5 — BESI_enrichi mis a jour dans master_dataset.csv",
        "MODULE 6 — Figure dual-axis NLP vs IPC (300 DPI) sauvegardee",
        "",
        "BESI_enrichi = 0.35×Trends + 0.25×NLP_Maroc + 0.20×YouTube + 0.10×Reddit + 0.10×|dIPC|",
    ],
    num_label="11"
)

# =============================================================================
# DIAPOSITIVE 13 — VERDICT H1
# =============================================================================
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 1.05, C_NAVY)
add_rect(sl, 0, 0, 0.2, 1.05, C_ORANGE)
add_text(sl, "Verdict H1 & Conclusion", 0.35, 0.1, 12.5, 0.85,
         font_size=22, bold=True, color=C_WHITE)
add_rect(sl, 0, 1.05, 13.33, 6.45, C_GREY)

# Verdict box
add_rect(sl, 0.4, 1.2, 12.5, 0.55, C_ORANGE)
add_text(sl, "H1 : PARTIELLEMENT REJETEE — Lead time reel = 12 mois (non 1-2 mois)",
         0.5, 1.25, 12.3, 0.45,
         font_size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# Colonne gauche : REJETE
add_rect(sl, 0.4, 1.85, 5.9, 2.4, C_RED)
add_text(sl, "CE QUI EST REJETE", 0.5, 1.9, 5.7, 0.4,
         font_size=13, bold=True, color=C_WHITE)
add_text(sl,
    "✗  Lead time 1-2 mois\n   => Lead time reel : 12 mois\n\n"
    "✗  SARIMAX+BESI ameliore le RMSE\n   => SARIMA reste meilleur\n\n"
    "✗  Relation lineaire stable\n   => Rupture 2022 confirmee",
    0.5, 2.35, 5.7, 1.8, font_size=12, color=C_WHITE)

# Colonne droite : CONFIRME
add_rect(sl, 6.9, 1.85, 5.9, 2.4, C_GREEN)
add_text(sl, "CE QUI EST CONFIRME", 7.0, 1.9, 5.7, 0.4,
         font_size=13, bold=True, color=C_WHITE)
add_text(sl,
    "✓  BESI precede l'IPC de 12 mois\n   (correlation croisee significative)\n\n"
    "✓  Granger : 7/7 features p < 0.05\n\n"
    "✓  BESI_trends : meilleur AIC -545\n   (information statistique reelle)\n\n"
    "✓  Early Warning : Rappel=100%",
    7.0, 2.35, 5.7, 1.8, font_size=12, color=C_WHITE)

# Implication
add_rect(sl, 0.4, 4.35, 12.5, 1.1, RGBColor(0xE8, 0xF4, 0xF8))
add_text(sl, "IMPLICATION PRATIQUE",
         0.55, 4.4, 12.2, 0.35, font_size=12, bold=True, color=C_BLUE)
add_text(sl,
    "L'horizon de BESI est MACROECONOMIQUE (12 mois), non tactique (1-2 mois).\n"
    "Utile pour les politiques contracycliques (credit, subventions alimentaires) avec un horizon de 12 mois,\n"
    "completant les analyses statistiques traditionnelles du HCP.",
    0.55, 4.75, 12.2, 0.65, font_size=12, color=C_DARK)

# Chiffres cles
add_text(sl,
    "CHIFFRES CLES :   180 obs  |  SARIMA RMSE=0.00272  |  BESI_trends AIC=-545.88  |  Lead=12 mois  |  Rappel=100%  |  F1=0.82",
    0.3, 5.6, 12.7, 0.45,
    font_size=12, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)

add_text(sl, "12", 6.1, 5.85, 1.1, 1.1, font_size=48, bold=True,
         color=C_ORANGE, align=PP_ALIGN.CENTER)
add_text(sl, "mois d'alerte\nprecoce avant\nle choc 2022",
         5.0, 6.1, 3.3, 1.1, font_size=11, color=C_NAVY, align=PP_ALIGN.CENTER)

# =============================================================================
# DIAPOSITIVE 14 — PHRASE DE POSITIONNEMENT
# =============================================================================
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, C_NAVY)
add_rect(sl, 0, 0, 0.25, 7.5, C_ORANGE)
add_rect(sl, 6.5, 2.2, 6.6, 3.1, RGBColor(0x25, 0x40, 0x60))

add_text(sl, "PHRASE DE POSITIONNEMENT — ORAL", 0.5, 0.4, 12.3, 0.6,
         font_size=16, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)

quote = (
    "\"Je reste dans le cadre SARIMA/SARIMAX du cours, mais j'introduis\n"
    "une dimension comportementale multi-sources pour tester la stabilite\n"
    "structurelle apres 2022 et quantifier la capacite d'alerte precoce\n"
    "des signaux digitaux.\n\n"
    "Mes resultats montrent que SARIMA(2,1,1)x(0,1,1)[12] reste le meilleur\n"
    "modele en walk-forward (RMSE=0.00272), mais BESI_trends ameliore\n"
    "l'AIC de -502 a -546 — confirmant une information statistique reelle.\n\n"
    "BESI detecte le stress economique 12 mois d'avance — offrant un signal\n"
    "macroeconomique robuste pour les politiques contracycliques au Maroc.\""
)
add_text(sl, quote, 0.5, 1.1, 5.8, 5.8,
         font_size=14, color=C_WHITE, italic=True)

add_text(sl, "RESULTATS EN UN COUP D'OEIL", 6.6, 2.3, 6.2, 0.45,
         font_size=13, bold=True, color=C_ORANGE)
kpis = [
    ("0.00272", "RMSE SARIMA"),
    ("-545.88", "AIC SARIMAX_BT"),
    ("12 mois", "Lead time BESI"),
    ("100%",    "Rappel Early Warning"),
    ("0.82",    "F1-Score alerte"),
    ("90.9%",   "Persistance High Stress"),
]
for i, (val, lbl) in enumerate(kpis):
    row, col = divmod(i, 2)
    x = 6.6 + col * 3.1
    y = 2.85 + row * 0.95
    add_rect(sl, x, y, 2.85, 0.82, RGBColor(0x1F, 0x77, 0xB4))
    add_text(sl, val, x, y+0.04, 2.85, 0.42,
             font_size=18, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)
    add_text(sl, lbl, x, y+0.45, 2.85, 0.32,
             font_size=10, color=C_WHITE, align=PP_ALIGN.CENTER)

add_text(sl, "Douae Ahadji & Adama Basse  —  ENSAM Meknes  —  Séries Temporelles  —  Mai 2026",
         0.5, 7.0, 12.3, 0.4,
         font_size=12, color=RGBColor(0x88,0x99,0xAA),
         align=PP_ALIGN.CENTER, italic=True)

# =============================================================================
# SAUVEGARDE
# =============================================================================
out_path = ROOT / "BESI_Presentation_Soutenance.pptx"
prs.save(str(out_path))
print(f"PowerPoint sauvegarde : {out_path}")
print(f"Nombre de diapositives : {len(prs.slides)}")
